from pptx import Presentation
from pptx.util import Inches,Cm,Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def ppt_template(student_name,image):
    '''
    :param student_name:
    :param image:
    :return:
    '''
    ppt = Presentation()

    slide1 = ppt.slides.add_slide(ppt.slide_layouts[6])#标准的ppt 布局
    #第一个shape 名字
    left = Pt(280)
    top = Cm(1)
    width = Cm(5)
    height = Cm(2)
    new_text_shape_type = slide1.shapes.add_textbox(left,top,width,height)
    text_frame = new_text_shape_type.text_frame
    p = text_frame.add_paragraph()
    p.text = student_name
    p.alignment = PP_ALIGN.CENTER
    p.font.italic = True
    p.font.size=Pt(20)
    p.font.color.rgb = RGBColor(225, 0, 0)


    #第三个shape 图片
    new_text_shape_image = slide1.shapes.add_picture(image_file=image,
                                                     #width=Inches(3),
                                                     height=Pt(300), #注意这里限制一个边另一个边等比例
                                                     left=Cm(6),
                                                     top=Cm(5))
    ppt.save('ppts/%s.pptx'%student_name)





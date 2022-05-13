from pptx import Presentation
from pptx.util import Inches,Cm,Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from openpyxl import load_workbook

from pptx.enum.text import MSO_VERTICAL_ANCHOR

import os

def ppt_dealer(type,title,content,image,ppt):
    #一个slide 包含三个shapes ppt = Pres()
    slide1 = ppt.slides.add_slide(ppt.slide_layouts[6])#标准的ppt 布局
    #第一个shape 类别
    left = Cm(1)
    top = Cm(1)
    width = Cm(5)
    height = Cm(2)
    new_text_shape_type = slide1.shapes.add_textbox(left,top,width,height)
    text_frame = new_text_shape_type.text_frame
    p = text_frame.add_paragraph()
    p.text = type #'布艺沙发'
    p.alignment = PP_ALIGN.CENTER
    p.font.italic = True
    p.font.size=Pt(20)
    p.font.color.rgb = RGBColor(225, 0, 0)

    #第二个shape 家具的名字
    left = Cm(7.75)
    top = Cm(2)
    width = Cm(10)
    height = Cm(5)
    new_text_shape_title = slide1.shapes.add_textbox(left,top,width,height)
    text_frame = new_text_shape_title.text_frame
    p = text_frame.add_paragraph()
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    p.font.size=Pt(15)

    #第三个shape 图片
    new_text_shape_image = slide1.shapes.add_picture(image_file=image,
                                                     #width=Inches(3),
                                                     height=Pt(200), #注意这里限制一个边另一个边等比例
                                                     left=Cm(6),
                                                     top=Cm(5))

    #第四个shape 家具描述
    left = Cm(5)
    top = Cm(15)
    width = Cm(15)
    height = Cm(5)
    new_text_shape_content = slide1.shapes.add_textbox(left,top,width,height)
    text_frame = new_text_shape_content.text_frame
    text_frame.word_wrap = True #设定 当前的frame为自动换行
    p = text_frame.add_paragraph() #添加段落
    p.text = content
    p.font.size=Pt(12)



def main_dealer():
    ppt = Presentation() #创建ppt对象

    wb = load_workbook('家具素材.xlsx')
    ws = wb.worksheets[0]
    for r in list(ws.iter_rows())[1:]:

        row = [data.value for data in r]  # 注意这个推到列表非常重要

        if row[0] == '布艺沙发':
            ppt_dealer(type=row[0],title=row[1].split('-')[0],
                       content=row[2],image='布艺沙发/'+row[1]+'.jpeg',ppt=ppt)

        elif row[0] == '床':
            ppt_dealer(type=row[0], title=row[1].split('-')[0],
                       content=row[2], image='床/' + row[1] + '.jpeg', ppt=ppt)

        elif row[0] == '餐桌':
            ppt_dealer(type=row[0], title=row[1].split('-')[0],
                       content=row[2], image='餐桌/' + row[1] + '.jpeg', ppt=ppt)

    ppt.save('myFurniture.pptx')

if __name__ == '__main__':
    main_dealer()


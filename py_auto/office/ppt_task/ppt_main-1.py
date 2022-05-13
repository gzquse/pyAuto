from pptx import Presentation
from pptx.util import Inches, Cm, Pt, Mm
from pptx.enum.text import PP_ALIGN

ppt = Presentation()  # 创建ppt对象

# #创建 一个 slide
# slide1 = ppt.slides.add_slide(ppt.slide_layouts[1])#标准的ppt 布局
# #在当前slide获取所有的shape 框子
# shape = slide1.shapes.placeholders # 这个是获取你现在布局内的已经有的“框子”
#                     #注意 placeholder就是占位的意思
#                     #其实就是在定位你要修改的框子
#                     #然后这个就跟字典赋值一样了
# shape[0].text = '第一个内容'
# shape[1].text = '第二个内容'
#
# #在shape里添加段落 可以设置段落如文字的属性 大小
# paragraph = shape[1].text_frame.add_paragraph()
# paragraph.text = '新文字内容'
# paragraph.font.size = Pt(30)
# paragraph.font.name = '宋体'
#

# 以上都是使用现有的shape ，我们需要自己设置一个 shape

#
slide1 = ppt.slides.add_slide(ppt.slide_layouts[6])  # 标准的ppt 布局

new_text_shape = slide1.shapes.add_textbox(left=Cm(10), top=Cm(10), width=Cm(10), height=Cm(10))
text_frame = new_text_shape.text_frame
p = text_frame.add_paragraph()

p.text = "我自己的文本1111"
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(10)

ppt.save('aaa.pptx')  # 保存
